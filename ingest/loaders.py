import io
import os
import tempfile
from pathlib import Path
from typing import Optional, Tuple
from utils.file_utils import ext
from utils.ocr_utils import ocr_image

# Text file loaders
def load_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def load_md(path: str) -> str:
    import markdown2
    raw = Path(path).read_text(encoding="utf-8", errors="ignore")
    text = markdown2.markdown(raw)
    # strip HTML tags quickly
    return _html_to_text(text)

def load_pdf(path: str) -> str:
    from pdfminer.high_level import extract_text
    return extract_text(path) or ""

def load_docx(path: str) -> str:
    import docx
    d = docx.Document(path)
    return "\n".join(p.text for p in d.paragraphs)

def load_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

# Images → OCR
def load_image(path: str) -> str:
    return ocr_image(path)

# Audio/Video → transcription (Whisper local)
def load_media(path: str) -> str:
    from faster_whisper import WhisperModel
    import ffmpeg

    model = WhisperModel(model_size_or_path="large-v3" if _is_long(path) else "small")
    segments, _ = model.transcribe(path, beam_size=5)
    return "\n".join(s.text for s in segments if s.text).strip()

def _is_long(path: str) -> bool:
    # crude duration check via ffprobe
    import json, subprocess
    cmd = ["ffprobe", "-v", "quiet", "-print_format", "json", "-show_format", path]
    out = subprocess.check_output(cmd).decode()
    duration = float(json.loads(out)["format"].get("duration", 0.0))
    return duration > 1200  # >20 minutes → use larger model

# YouTube → captions or audio download + transcription
def load_youtube(url: str) -> str:
    from youtube_transcript_api import YouTubeTranscriptApi, TranscriptsDisabled, NoTranscriptFound
    vid = _extract_video_id(url)
    if not vid:
        return ""
    try:
        transcript = YouTubeTranscriptApi.get_transcript(vid, languages=["en"])
        return "\n".join([x["text"] for x in transcript])
    except (TranscriptsDisabled, NoTranscriptFound):
        # fallback: download audio and transcribe
        import yt_dlp, tempfile
        with tempfile.TemporaryDirectory() as td:
            outp = os.path.join(td, "%(id)s.%(ext)s")
            ydl_opts = {"format": "bestaudio/best", "outtmpl": outp, "quiet": True}
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                info = ydl.extract_info(url, download=True)
                audio_path = ydl.prepare_filename(info)
            return load_media(audio_path)

def _html_to_text(html: str) -> str:
    # minimal HTML stripper
    import re
    return re.sub("<[^<]+?>", "", html).strip()

def _extract_video_id(url: str) -> Optional[str]:
    import re
    # handles youtu.be and youtube.com
    patterns = [
        r"youtu\.be/([A-Za-z0-9_-]{6,})",
        r"v=([A-Za-z0-9_-]{6,})"
    ]
    for p in patterns:
        m = re.search(p, url)
        if m:
            return m.group(1)
    return None

# Router
def load_any(path_or_url: str) -> Tuple[str, str]:
    """
    Returns (text, source_label). Accepts file path or YouTube URL.
    """
    if path_or_url.startswith("http"):
        return load_youtube(path_or_url), f"YouTube::{path_or_url}"

    e = ext(path_or_url)
    if e in [".txt"]:
        return load_txt(path_or_url), path_or_url
    if e in [".md"]:
        return load_md(path_or_url), path_or_url
    if e in [".pdf"]:
        return load_pdf(path_or_url), path_or_url
    if e in [".docx"]:
        return load_docx(path_or_url), path_or_url
    if e in [".pptx"]:
        return load_pptx(path_or_url), path_or_url
    if e in [".png", ".jpg", ".jpeg"]:
        return load_image(path_or_url), path_or_url
    if e in [".mp3", ".mp4", ".wav", ".m4a"]:
        return load_media(path_or_url), path_or_url

    return "", path_or_url
